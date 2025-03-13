#!/usr/bin/env python3

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT, MSO_VERTICAL_ANCHOR, MSO_UNDERLINE
from itertools import takewhile
from create_images_from_text import *
import re

green_color = RGBColor(0, 200, 0)
black_color = RGBColor(50, 50, 50)
white_color = RGBColor(245, 245, 245)
image_dir = 'images'

# Define patterns for bold, italic, and underline
patterns = [
    (r"\*(.*?)\*", {"bold": True}),  # Bold: *text*
    (r"_(.*?)_", {"italic": True}),  # Italic: _text_
    (r"__(.*?)__", {"underline": MSO_UNDERLINE.SINGLE_LINE}),  # Underline: __text__
    (r"[1-3]*[ ]*[Song of]*[A-Za-z]*[ ][1-9]+:[0-9-,]*", {"bold": True})
]


class Spec:
    width = None
    height = None
    margin = None
    margin_left = None
    margin_top = None

    title_font = None
    content_font = None
    font_spacing = None
    font_name = 'Arial'

    title_height = None
    content_position = None
    content_height_1 = None
    content_height_2 = None
    textbox_width = None

    def set_ppt_spec_for_live(self):
        # footer only
        self.width = 1900
        self.height = 300
        self.margin = 20
        self.margin_left = self.margin
        self.margin_top = self.margin

        self.title_font = 36
        self.content_font = 28
        self.font_spacing = 1.2
        self.font_name = 'Arial'

        self.update_dimensions()

    def set_ppt_spec_for_main(self):
        # 16:9 ratio
        self.width = 1920
        self.height = 1080
        self.margin = 50

        self.title_font = 70
        self.content_font = 55
        self.font_spacing = 1.5
        self.font_name = 'Arial'

        self.update_dimensions()

    def set_image_spec_for_live(self):
        # footer only
        self.width = 1900
        self.height = 300
        self.margin = 20
        self.margin_left = self.margin * 3
        self.margin_top = self.margin

        self.title_font = 40
        self.content_font = 32
        self.font_spacing = 15
        self.font_name = '/System/Library/Fonts/Supplemental/Arial.ttf'

        self.update_dimensions()

    def update_dimensions(self):
        self.title_height = self.margin + self.title_font
        self.content_position = self.margin + self.title_height
        self.content_height_1 = self.height - 2 * self.margin
        self.content_height_2 = self.height - 2 * self.margin - self.title_height
        self.textbox_width = self.width - 2 * self.margin


def create_presentation(width, height):
    presentation = Presentation()
    presentation.slide_width = Pt(width)
    presentation.slide_height = Pt(height)
    return presentation


def new_slide(presentation):
    # 0. Title (presentation title slide)
    # 1. Title and Content
    # 2. Section Header (sometimes called Segue)
    # 3. Two Content (side by side bullet textboxes)
    # 4. Comparison (same but additional title for each side by side content box)
    # 5. Title Only
    # 6. Blank
    # 7. Content with Caption
    # 8. Picture with Caption
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])

    # Add bg color
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = black_color

    # Add bg image
    # slide.shapes.add_picture(background_image_path, Inches(0), Inches(0), Inches(9), Inches(6))

    return slide


def is_title(line):
    return line.startswith('#')


def get_indent(line):
    n = count_leading_spaces(line)
    if n % 2 != 0:
        raise RuntimeError(f'Remove extra space at "{line}"')
    return int(n / 2)


def count_leading_spaces(s):
    return sum(1 for _ in takewhile(lambda x: x == " ", s))


def set_slide_title(textFrame, titleText):
    # add_formatted_text_runs(textFrame.paragraphs[0], titleText, title_font)
    textFrame.paragraphs[0].text = titleText


def add_bullets(text, run):
    text = text.strip()
    if text.startswith('*') or text.startswith('-'):
        text = text.replace('*', '● ', 1)
        # text = text.replace('-', '- ', 1)
    return text


def add_formatted_text_runs(para, text, spec):
    # Track the last position in the text
    last_pos = 0

    # Iterate through patterns and apply formatting
    for match in re.finditer(r"|".join([p[0] for p in patterns]), text):
        start, end = match.span()
        if start > last_pos:
            # Add unformatted text
            run = para.add_run()
            run.text = add_bullets(text[last_pos:start], run)
            # run.font.size = Pt(font_size)

        # Identify matched formatting
        for pattern, styles in patterns:
            if re.fullmatch(pattern, match.group()):
                run = para.add_run()
                # text_value = next(ftext for ftext in match.groups() if ftext is not None)
                text_value = match.group()
                for ftext in match.groups():
                    if ftext:
                        text_value = ftext
                        break
                run.text = add_bullets(text_value, run)
                for style, value in styles.items():
                    setattr(run.font, style, value)
                break
        last_pos = end

    # Add remaining unformatted text
    if last_pos < len(text):
        run = para.add_run()
        run.text = add_bullets(text[last_pos:], run)
        run.font.size = Pt(spec.content_font)


def add_slide_content(para, text, indent, spec):
    format_para_for_content(para, indent, spec)

    # text = text.strip()
    if text.startswith('*') or text.startswith('-'):
        para.bullet = True
        #     text = text.replace('*', '● ', 1)
    #     # text = text.replace('-', '- ', 1)
    else:
        para.bullet = False

    # Ensure it's not indented as a bullet
    # Level determines bullet/numbering style (0 is top level)
    para.level = indent
    para.space_before = Pt(0)  # Adjust spacing before the paragraph
    para.space_after = Pt(0)  # Adjust spacing after the paragraph
    para.left_margin = Pt(36)  # Indents the first line
    para.indent = Pt(-18)  # Hanging indent for wrapped text

    add_formatted_text_runs(para, text, spec)
    # para.text = text

    # TODO: Doesn't work
    # para._element.get_or_add_pPr().set("numId", "1")  # Enable numbering
    # bullet = para._element.get_or_add_pPr().get_or_add_buChar()
    # bullet.set("char", bullet_char)


def create_title_textbox(slide, spec):
    textBox = slide.shapes.add_textbox(Pt(spec.margin), Pt(spec.margin), Pt(spec.textbox_width), Pt(spec.title_height))
    textFrame = textBox.text_frame
    textFrame.word_wrap = True
    textFrame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    para = textFrame.paragraphs[0]
    format_para_for_title(para, spec)
    return textFrame


def create_content_textbox(slide, has_title, spec):
    # slide.shapes.placeholders[1].text_frame
    if has_title:
        textBox = slide.shapes.add_textbox(Pt(spec.margin), Pt(spec.content_position), Pt(spec.textbox_width),
                                           Pt(spec.content_height_1))
    else:
        textBox = slide.shapes.add_textbox(Pt(spec.margin), Pt(spec.margin), Pt(spec.textbox_width),
                                           Pt(spec.content_height_2))

    textFrame = textBox.text_frame
    textFrame.clear()
    textFrame.word_wrap = True
    textFrame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    return textFrame


def format_para_for_title(para, spec):
    para.line_spacing = spec.font_spacing
    para.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
    para.line_spacing = spec.font_spacing
    para.font.name = spec.font_name
    para.font.size = Pt(spec.title_font)
    para.font.bold = True
    para.font.color.rgb = white_color


def format_para_for_content(para, indent, spec):
    para.line_spacing = spec.font_spacing
    para.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
    para.font.name = spec.font_name
    para.font.size = Pt(spec.content_font - indent * 2)
    para.font.bold = False
    para.font.color.rgb = white_color


def convert_text_to_presentation(slides, spec, name):
    # Create presentation
    presentation = create_presentation(spec.width, spec.height)

    # Add slides to presentation
    for slide_content in slides:
        slide = new_slide(presentation)
        has_title = is_title(slide_content[0])
        title_textFrame = create_title_textbox(slide, spec)
        content_textFrame = create_content_textbox(slide, has_title, spec)
        content_textFrame.clear()
        first_content_line = True

        for line in slide_content:
            if is_title(line):
                titleText = line.strip('# ')
                set_slide_title(title_textFrame, titleText)
            elif first_content_line:
                first_content_line = False
                indent = get_indent(line)
                content_para = content_textFrame.add_paragraph()
                # content_para = content_textFrame.paragraphs[0]
                add_slide_content(content_para, line.strip(), indent, spec)
            else:
                indent = get_indent(line)
                content_para = content_textFrame.add_paragraph()
                add_slide_content(content_para, line.strip(), indent, spec)

    presentation.save(name)


def get_slides_list_from_text(filename):
    with open(filename, 'r') as contentfile:
        content = contentfile.readlines()

    # Read ppt content into a slide list
    slides = []
    slide = []
    for line in content:
        line = line.rstrip('\n')
        if line == '==':
            # end of slide, create new slide
            slides.append(slide)
            slide = []
            continue
        slide.append(line)

    # Debug slide list
    for s in slides:
        print('\n'.join(s))

    return slides


def main():
    slides = get_slides_list_from_text('content.txt')

    main_spec = Spec()
    main_spec.set_ppt_spec_for_main()
    convert_text_to_presentation(slides, main_spec, 'Presentation Main.pptx')

    live_spec = Spec()
    live_spec.set_ppt_spec_for_live()
    convert_text_to_presentation(slides, live_spec, 'Presentation LIVE.pptx')

    # live_spec.set_image_spec_for_live()
    # init_images_dir(image_dir)
    # create_ppt_images(slides, live_spec, image_dir)

    print('\nDone.')


if __name__ == "__main__":
    main()
