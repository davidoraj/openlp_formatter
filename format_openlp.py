import re
import os
import string
from collections import OrderedDict
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT, MSO_VERTICAL_ANCHOR

from create_lyrics_images import *

# import xml.dom.minidom

# README
# Max chars per line = 38

default_lines_per_slide = 2
section_lookup = {'V': 'Verse', 'C': 'Chorus', 'P': 'Pre-chorus', 'B': 'Bridge', 'T': 'Other', 'O': 'Other'}
sections = {'Verse', 'Chorus', 'Pre-chorus', 'Bridge', 'Tag', 'Instrumental', 'BREAK'}
ignore_section = {'BREAK'}
sections_regex = '|'.join([f'^{section}[0-9 ]*:' for section in sections])
repeat_section_regex = '|'.join([f'^{section}[0-9 ]*$' for section in sections])

# Variables used for creating song xml
verse_parts = list(string.ascii_lowercase)  # Creates ['a', 'b', 'c', ..., 'z']
line_delim = '<br/>'
section_delim = '<br/><br/>'

lyrics_ppt_file_green = "Lyrics GREEN.pptx"
lyrics_ppt_file_main = "ICC Worship Lyrics.pptx"
background_image_path = 'backgrounds/ICC_slides_template.001.jpeg'

margin = 0.2
total_width = 8
total_height = 6
inches_margin_left = Inches(margin)
inches_margin_right = Inches(margin)
inches_margin_top = Inches(margin)
inches_total_width = Inches(total_width)
inches_total_height = Inches(total_height)
inches_total_width_full = Inches(total_width - margin * 2)
inches_total_height_full = Inches(total_height - margin * 2)
inches_total_width_half = Inches(total_width / 2)
inches_total_height_half = Inches(total_height / 2)

font_name = 'Gautami'

# Fonts for green slides (Live)
font_size_title_green = 26
font_size_green = 14  # 16
font_spacing_green = 22  # 24

# Fonts for main slides
font_size_title_main = 38
font_size_main = 26
font_spacing_main = 38

# c-chorus, n-verse, p-prechorus, b-bridge, c2

song_text = ""
with open('lyrics_text/lyrics_02-27-22.txt', 'r') as lyricsfile:
    song_text = lyricsfile.readlines()


# Checks if a given line is the definition of a section. Below are a few examples
# Chorus:
# Verse 2:
# Verse 1:[3] <-- specifies section-specific lines_per_slide
def is_section(line):
    if re.match(sections_regex, line):
        colon_splits = line.split(":")
        lps = None
        # Check if section-specific lines_per_slide is given [n]
        if len(colon_splits[1]) > 0:
            # Supports numbers 1-9 only
            lps = int(colon_splits[1][1])

        words = colon_splits[0].split()

        if len(words) > 1:
            return words[0], words[1], lps
        else:
            return words[0], 1, lps
    else:
        return None, None, None


def is_repeat_Section(line):
    section = None
    id = 1
    if re.match(repeat_section_regex, line):
        words = line.split()
        section = words[0]
        if len(words) > 1:
            if re.match('[0-9]+', words[1]):
                id = words[1]
            else:
                id = '1'

    return section, id


# Parses the song from given annotated English and Telugu text
def get_song_annotated(one, two, lines_per_slide, title):
    sections = dict()
    section_lines = dict()
    order = []

    for line in one:
        if not line:
            continue

        section, id, sec_lines = is_section(line)
        rsection, rid = is_repeat_Section(line)
        if section:
            if section[0] == 'I':  # Instrumental
                sid = 'O' + str(id)
            else:
                sid = section[0] + str(id)
            if section not in ignore_section:
                order.append(sid)

            # If section-specific lines per slide is given use it
            section_lines[sid] = sec_lines or lines_per_slide

        elif rsection:
            sid = rsection[0] + str(rid)
            if rsection not in ignore_section:
                order.append(sid)
            sid = None
        else:
            try:
                if not sid:
                    print('Missing annotation or title at "{}"'.format(line))
                    exit(1)
                else:
                    key = sid + '_one'
                    content = sections.get(key, [])
                    content.append(line)
                    sections[key] = content
            except UnboundLocalError:
                print('Missing annotation or title at "{}"'.format(line))
                exit(1)

    section_map = set()
    i = 0
    for o in order:
        key = o + '_two'
        if key in section_map:
            continue

        while i < len(two):
            line = two[i]

            if line:
                content = sections.get(key, [])
                content.append(line)
                sections[key] = content
                i = i + 1
            else:
                section_map.add(key)
                i = i + 1
                break

    return {'sections': sections,
            'section_lines': section_lines,
            'order': order,
            'verse_order': ' '.join(order),
            'lines_per_slide': lines_per_slide,
            'title': title}


# Parses the songs from text (list of lines)
def get_song_objects(song_lines_list):
    songs = []

    for song_lines in song_lines_list:
        line = song_lines[0][0]  # first line ---
        title = song_lines[0][1].title()  # second line
        one = song_lines[0][2:]
        two = song_lines[1]

        # Get lines per slide or default, (for next song)
        if re.match("---[0-9]+", line):
            lines_per_slide = int(line[3:])
        else:
            lines_per_slide = default_lines_per_slide

        # Parse song to object (based on annotations)
        song = get_song_annotated(one, two, lines_per_slide, title)

        # Create slides (English + Telugu combined)
        create_song_slide_deck(song)

        songs.append(song)

    return songs


def append_section_lyrics_xml(id, verse_part_counter, slides):
    counter = verse_part_counter.get(id, -1) + 1
    verse_part_counter[id] = counter
    vid = id.lower() + verse_parts[counter]
    return vid, slides[1:]


def append_section_slides(label, id, deck, section, lines_per_slide=default_lines_per_slide):
    # Init vars
    n = len(section)
    line_count = n if n < lines_per_slide else lines_per_slide
    i_last_slide = n if n % lines_per_slide == 0 else n - (n % lines_per_slide)  # index for last slide

    verse_part_counter = {}  # verse index (with alphabet counter, v1a, v1b etc)
    section_lyrics_dict = OrderedDict()

    # list of lines in a slide (added to given deck)

    # All but last slide (last n % lines_per_slide lines in the section)
    for i in range(0, i_last_slide, line_count):
        slides = []
        slides.append(label)
        for j in range(line_count):
            slides.append(section[j + i])

        vid, slide_lines = append_section_lyrics_xml(id, verse_part_counter, slides)
        section_lyrics_dict[vid] = slide_lines
        deck.append(slides)

    # Last slide for the section
    if i_last_slide < n:
        slides = []
        slides.append(label)
        for i in range(i_last_slide, n):
            slides.append(section[i])

        vid, slide_lines = append_section_lyrics_xml(id, verse_part_counter, slides)
        section_lyrics_dict[vid] = slide_lines
        deck.append(slides)

    return section_lyrics_dict


def merge_section_slides(main_deck, temp_deck, i):
    line_count = len(temp_deck[i])

    # for each line in a section
    for j in range(1, line_count):
        main_deck.append(temp_deck[i][j])


def merge_decks(deck1, deck2):
    # deck1 and deck2 lengths should be the same at this point
    deck = []
    n = len(deck1)

    # for each section
    for i in range(n):
        # add section label
        deck.append(deck1[i][0])
        merge_section_slides(deck, deck1, i)
        if deck2:
            deck.append("")  # empty line
            merge_section_slides(deck, deck2, i)

    return deck


# Takes two lists of tuples, and returns an OrderedDict of named lyrics (c1a->..., c1b->..., v1a->... etc.)
def merge_lyrics_xml(dict1, dict2):
    lyrics_xml = OrderedDict()

    for id in dict1.keys():
        xml_string = line_delim.join(dict1[id])
        if dict2:
            # (Telugu on top, English at the bottom)
            xml_string = line_delim.join(dict2[id]) + section_delim + xml_string
        lyrics_xml[id] = xml_string

    return lyrics_xml


# Given an annotated song (sections/order)
# it returns the slides for the song (English & Telugu)
def create_song_slide_deck(song):
    section_map = set()
    deck1 = []  # english
    deck2 = []  # telugu
    lyrics1_xml_dict = OrderedDict()
    lyrics2_xml_dict = OrderedDict()

    for o in song['order']:
        one = o + '_one'
        two = o + '_two'
        if one in song['sections']:
            section1 = song['sections'][one]
        if two in song['sections']:
            section2 = song['sections'][two]
        else:
            section2 = None

        n = len(section1)

        # Validation
        if section2 and len(section1) != len(section2):
            raise SyntaxError("Lines don't match for song: " + str(song))

        label = '---[{sec}:{id}]---'.format(sec=section_lookup[o[0]], id=o[1])

        # skip duplicate section
        if label in section_map:
            continue
        else:
            section_map.add(label)

        # Get lines_per_slide for given section
        # If not specified for a section, it defaults to song lines_per_slide
        # If not then, it defaults to global (default_lines_per_slide)
        section_lines_per_slide = song['section_lines'][o]

        # Convert a section (with n lines) to k slides, and put in a dict with vid (verse id) as key
        lyrics1_xml_dict.update(
            append_section_slides(label, o, deck1, section1, section_lines_per_slide))
        if section2:
            lyrics2_xml_dict.update(
                append_section_slides(label, o, deck2, section2, section_lines_per_slide))

    # Merge decks into a a single deck (song text in "edit-all")
    song['deck'] = merge_decks(deck1, deck2)

    # Merge lyrics xml for song import
    song['lyrics_xml'] = merge_lyrics_xml(lyrics1_xml_dict, lyrics2_xml_dict)

    return


# Returns list of tuples (englist_lines_list, optional_telugu_lines_list)
def get_song_lines(lines):
    song_lines_list = []
    one = []
    two = []
    curlist = None

    for i in range(len(lines)):
        line = lines[i]
        line = line.strip().strip("\n")

        # TODO: Capitalize first alphabet ignore numbers
        if re.match('^[a-z]', line):
            line = line.capitalize()

        if line.startswith("---"):
            if one:
                song_lines_list.append((one, two))
            one = []
            two = []
            curlist = one
            curlist.append(line)
        elif line == "--":
            curlist = two
        elif line:
            curlist.append(line)
        elif curlist is not None:
            curlist.append("")

    return song_lines_list


def get_xml(tag, text, props=''):
    return '<{tag} {props}>{text}</{tag}>'.format(tag=tag, text=text, props=props)


def save_to_xml(song):
    start = """<?xml version='1.0' encoding='UTF-8'?>
<song xmlns="http://openlyrics.info/namespace/2009/song" version="0.8" createdIn="OpenLP 2.4.6" modifiedIn="OpenLP 2.4.6">"""

    properties_text = \
        get_xml('titles', get_xml('title', song['title'])) + \
        get_xml('authors', get_xml('author', 'Unknown')) + \
        get_xml('verseOrder', f"i1 {song['verse_order'].lower()} i2")
    properties = '\n' + get_xml('properties', properties_text) + '\n'

    # Create a list of verses with name=verse_id
    # song['lyrics_xml'] is an ordered dict with key=verse_id, value=[slide-lyrics]
    lyrics_xml_list = []

    # Add title
    lyrics_xml_list.append(
        get_xml('verse', get_xml('lines', song['title']), 'name="i1"'))

    for vid in song['lyrics_xml']:
        verse_lines = get_xml('lines', ''.join(song['lyrics_xml'][vid]))
        lyrics_xml_list.append(get_xml('verse', verse_lines, 'name="{}"'.format(vid)))
    end = "\n</song>"

    # Add song ending slide
    lyrics_xml_list.append(
        get_xml('verse', get_xml('lines', ''), 'name="i2"'))

    # Put combine xml content into one string
    song_xml = start + properties + get_xml('lyrics', '\n'.join(lyrics_xml_list)) + end

    # Write to file
    with open(f"songs_xml/{song['title']}.xml", 'w') as xml_file:
        # TODO: Find a way to pretty print without disturbing verse lines text
        # xml_file.write(xml.dom.minidom.parseString(song_xml).toprettyxml())
        xml_file.write(song_xml)
    return


def get_song_lyrics_content(song, i):
    lyrics_dict = song['lyrics_xml']  # odict_keys(['c1a', 'v1a', 'v2a', 'v3a', 'v4a', 'v5a', 'v6a'])
    j = 0
    two_langs = True if song['order'][0] + '_two' in song['sections'] else False
    content = []
    for id in song['order']:
        id = id.lower()
        j = j + 1
        for counter in verse_parts:
            vid = id + counter
            if vid in lyrics_dict:
                img_name = '{:0>2d}-{:0>2d}_{}.{}'.format(i, j, vid, img_format)
                if two_langs:
                    lyrics_two_langs = lyrics_dict[vid].split(line_delim + line_delim)
                    text1 = lyrics_two_langs[0].replace(line_delim, '\n')
                    text2 = lyrics_two_langs[1].replace(line_delim, '\n')
                    content.append((img_name, text1, text2))
                else:
                    text1 = lyrics_dict[vid].replace(line_delim, '\n')
                    content.append((img_name, text1, None))
            else:
                break

    return content


def get_green_slide(root):
    slide = root.slides.add_slide(root.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 200, 0)
    return slide


def get_blank_slide(root):
    slide = root.slides.add_slide(root.slide_layouts[6])
    slide.shapes.add_picture(background_image_path, Inches(0), Inches(0), Inches(9), Inches(6))
    return slide


def add_textbox_green(slide, text, left, top, width, height):
    textBox = slide.shapes.add_textbox(left, top, width, height)
    textFrame = textBox.text_frame
    para = textFrame.paragraphs[0]
    para.text = text
    para.line_spacing = Pt(font_spacing_green)
    para.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
    font = para.font
    font.name = font_name
    font.size = Pt(font_size_green)
    font.bold = True
    font.color.rgb = RGBColor(255, 255, 255)
    return


def add_textbox_main(slide, text, left, top, width, height):
    textBox = slide.shapes.add_textbox(left, top, width, height)
    textFrame = textBox.text_frame
    textFrame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    para = textFrame.paragraphs[0]
    para.text = text
    para.line_spacing = Pt(font_spacing_main)
    para.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
    font = para.font
    font.name = font_name
    font.size = Pt(font_size_main)
    font.bold = True
    font.color.rgb = RGBColor(255, 255, 255)
    return


def add_title_green(root, title):
    slide = get_green_slide(root)
    shapes = slide.shapes
    textBox = shapes.add_textbox(inches_margin_left,
                                 inches_margin_top,
                                 inches_total_width_full,
                                 inches_total_height_full)
    textFrame = textBox.text_frame
    para = textFrame.paragraphs[0]
    para.text = title
    para.line_spacing = Pt(font_spacing_green)
    para.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
    font = para.font
    font.name = font_name
    font.size = Pt(font_size_title_green)
    font.bold = True
    font.color.rgb = RGBColor(255, 255, 255)
    return


def add_title_main(root, title):
    slide = get_blank_slide(root)
    shapes = slide.shapes
    textBox = shapes.add_textbox(inches_margin_left,
                                 inches_margin_top,
                                 inches_total_width_full,
                                 inches_total_height_full)
    textFrame = textBox.text_frame
    textFrame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE  # Vertical align middle title
    para = textFrame.paragraphs[0]
    para.text = title
    para.line_spacing = Pt(font_spacing_main)
    para.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
    font = para.font
    font.name = font_name
    font.size = Pt(font_size_title_main)
    font.bold = True
    font.color.rgb = RGBColor(255, 255, 255)
    return


def save_to_ppt_green(content, root, title):
    try:
        os.remove(lyrics_ppt_file_green)
    except:
        print()

    add_title_green(root, "Worship")
    # add_title_green(root, '')  # BLANK slide
    add_title_green(root, title)

    for img_name, text1, text2 in content:
        slide = get_green_slide(root)

        if text2:
            add_textbox_green(slide, text1, inches_margin_left, inches_margin_top,
                              inches_total_width_half, inches_total_height_full)

            add_textbox_green(slide, text2, inches_total_width_half, inches_margin_top,
                              inches_total_width_half, inches_total_height_full)
        else:
            add_textbox_green(slide, text1, inches_margin_left, inches_margin_top,
                              inches_total_width_full, inches_total_height_full)

    return


def save_to_ppt_main(content, root, title):
    try:
        os.remove(lyrics_ppt_file_main)
    except:
        print()

    add_title_main(root, title)

    for img_name, text1, text2 in content:
        slide = get_blank_slide(root)

        if text2:
            add_textbox_main(slide, text1, inches_margin_left, inches_margin_top,
                             inches_total_width_full, inches_total_height_half)

            add_textbox_main(slide, text2, inches_margin_left, inches_total_height_half,
                             inches_total_width_full, inches_total_height_half)
        else:
            add_textbox_main(slide, text1, inches_margin_left, inches_margin_top,
                             inches_total_width_full, inches_total_height_full)

    return


def create_new_presentation():
    presentation = Presentation()
    presentation.slide_width = Inches(total_width)
    presentation.slide_height = Inches(total_height)
    return presentation


def main():
    # Split songs lines to a list
    song_lines_list = get_song_lines(song_text)

    # Parse all the songs (returns a list of song objects)
    songs = get_song_objects(song_lines_list)

    # Init empty presentation
    pptx_green = create_new_presentation()
    pptx_main = create_new_presentation()
    init_images_dir()

    i = 1
    # Create slides
    for song in songs:

        # Export to xml
        save_to_xml(song)

        # Export to images / ppt
        content = get_song_lyrics_content(song, i)
        # save_to_images(content)
        save_to_ppt_green(content, pptx_green, song['title'])
        save_to_ppt_main(content, pptx_main, song['title'])
        i = i + 1

        # Print text
        # print("----------------------------------")
        # print(song['verse_order'])
        # for line in song['deck']:
        #     print(line)

        print("----------------------------------")
        for id, text1, text2 in content:
            print('--')
            print(text1)
            if text2:
                print(text2)

    pptx_green.save(lyrics_ppt_file_green)
    pptx_main.save(lyrics_ppt_file_main)

    print("----------------------------------")
    i = 1
    for song in songs:
        print('{}. {}'.format(i, song['title']))
        i = i + 1


if __name__ == "__main__":
    main()
